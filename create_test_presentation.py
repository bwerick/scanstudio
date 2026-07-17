#!/usr/bin/env python3
"""
Create PowerPoint presentation about ScanStudio testing strategy
Focused on 3 verification signals: Duplicate-frame, Entity preservation, Retrieval
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point['text']
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = point['text']
        p.level = point.get('level', 0)
        p.font.size = Pt(point.get('size', 18))
    return slide

def add_code_slide(prs, title, code_text, description):
    """Add a slide with code example"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True

    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.italic = True

    # Code
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(5))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            code_frame.text = line
            p = code_frame.paragraphs[0]
        else:
            p = code_frame.add_paragraph()
            p.text = line
        p.font.name = 'Courier New'
        p.font.size = Pt(11)
    fill = code_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    return slide

def add_diagram_slide(prs, title, diagram_items):
    """Add a slide with visual diagram/flowchart"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Diagram content
    y_pos = 1.2
    for item in diagram_items:
        box = slide.shapes.add_textbox(
            Inches(item.get('x', 0.5)),
            Inches(y_pos),
            Inches(item.get('width', 9)),
            Inches(item.get('height', 0.8))
        )
        tf = box.text_frame
        tf.text = item['text']
        p = tf.paragraphs[0]
        p.font.size = Pt(item.get('size', 16))
        p.font.bold = item.get('bold', False)
        p.alignment = PP_ALIGN.CENTER if item.get('center', False) else PP_ALIGN.LEFT

        if item.get('bg_color'):
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*item['bg_color'])

        y_pos += item.get('height', 0.8) + 0.2

    return slide

# ==================== SLIDES ====================

# Slide 1: Title
add_title_slide(prs,
    "ScanStudio Testing Strategy",
    "Three Self-Verifiable Signals for Document Intelligence"
)

# Slide 2: The Problem
add_content_slide(prs, "The Challenge", [
    {'text': 'ScanStudio Pipeline: Video → PDF → Searchable Database'},
    {'text': '9 processing phases with complex algorithms', 'level': 1},
    {'text': 'How do we verify correctness without manual inspection?', 'size': 22},
    {'text': 'Traditional unit tests are insufficient', 'level': 1, 'size': 16},
    {'text': 'Generic computer vision metrics don\'t capture document intelligence goals', 'level': 1, 'size': 16},
    {'text': 'We need tests that verify end-to-end document understanding', 'size': 20},
])

# Slide 3: Three Verification Signals
add_content_slide(prs, "Three Self-Verifiable Signals", [
    {'text': 'A. Duplicate-Frame Agreement', 'size': 24},
    {'text': 'Same page captured multiple times → consistent results', 'level': 1, 'size': 16},
    {'text': 'B. Entity Preservation', 'size': 24},
    {'text': 'Critical entities (names, numbers, dates) survive processing', 'level': 1, 'size': 16},
    {'text': 'C. Retrieval Correctness', 'size': 24},
    {'text': 'Final database answers known queries accurately', 'level': 1, 'size': 16},
])

# Slide 4: Why These Three?
add_content_slide(prs, "Why These Verification Signals?", [
    {'text': '✅ Self-Verifiable'},
    {'text': 'No human judgment needed - automatic pass/fail', 'level': 1},
    {'text': '✅ End-to-End'},
    {'text': 'Test the entire pipeline, not isolated components', 'level': 1},
    {'text': '✅ Aligned with Real Use Cases'},
    {'text': 'Document intelligence goal: extract and retrieve information', 'level': 1},
    {'text': '✅ Connects to SIGMOD Research'},
    {'text': 'Ground truth maintenance for document lifecycle', 'level': 1},
])

# Slide 5: Signal A Overview
add_content_slide(prs, "Signal A: Duplicate-Frame Agreement", [
    {'text': 'Core Insight: Same page → Same output'},
    {'text': 'When scanning, pages are often captured 2-3 times'},
    {'text': 'Frame 150: Page 5 (first capture)', 'level': 1, 'size': 16},
    {'text': 'Frame 480: Page 5 (accidentally re-captured)', 'level': 1, 'size': 16},
    {'text': 'Test: Do both captures produce identical/similar results?', 'size': 20},
    {'text': 'High agreement → pipeline is stable and repeatable', 'level': 1},
    {'text': 'Low agreement → pipeline has inconsistency bugs', 'level': 1},
])

# Slide 6: Signal A - How It Works
add_diagram_slide(prs, "Duplicate-Frame Agreement: How It Works", [
    {'text': '1. Identify Duplicate Captures', 'size': 18, 'bold': True, 'bg_color': (230, 240, 255)},
    {'text': '   • Use perceptual hashing (pHash) to find near-identical frames', 'size': 14},
    {'text': '   • Hamming distance < 10 → same page', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '2. Process Both Captures Through Pipeline', 'size': 18, 'bold': True, 'bg_color': (230, 255, 230)},
    {'text': '   • Frame 150 → Crop → OCR → Text₁', 'size': 14},
    {'text': '   • Frame 480 → Crop → OCR → Text₂', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '3. Compute Agreement Score', 'size': 18, 'bold': True, 'bg_color': (255, 240, 230)},
    {'text': '   • Text similarity (Levenshtein distance)', 'size': 14},
    {'text': '   • Bounding box IoU', 'size': 14},
    {'text': '   • Entity overlap (names, dates, numbers)', 'size': 14},
])

# Slide 7: Signal A - Code Example
code_a = """def test_duplicate_frame_agreement():
    # Find duplicate page captures in test video
    duplicates = find_duplicate_frames(
        "test_video.mp4",
        phash_threshold=10
    )

    # duplicates = [(frame_150, frame_480), (frame_200, frame_520)]

    agreement_scores = []
    for frame_a, frame_b in duplicates:
        # Process both frames through full pipeline
        text_a = pipeline(extract_frame(frame_a))
        text_b = pipeline(extract_frame(frame_b))

        # Compute similarity
        similarity = levenshtein_similarity(text_a, text_b)
        agreement_scores.append(similarity)

    # Self-verifying: identical pages should have >95% agreement
    assert np.mean(agreement_scores) > 0.95
    ✅ AUTOMATIC VERIFICATION"""

add_code_slide(prs, "Signal A: Implementation",
    code_a,
    "Identical pages must produce consistent results"
)

# Slide 8: Signal A - Metrics
add_content_slide(prs, "Signal A: Agreement Metrics", [
    {'text': 'Text Similarity'},
    {'text': 'Levenshtein distance on OCR output', 'level': 1},
    {'text': 'Threshold: ≥95% agreement', 'level': 1},
    {'text': 'Bounding Box IoU'},
    {'text': 'Crop coordinates should be nearly identical', 'level': 1},
    {'text': 'Threshold: IoU ≥0.90', 'level': 1},
    {'text': 'Entity Overlap'},
    {'text': 'Extracted entities (names, dates) should match', 'level': 1},
    {'text': 'Threshold: ≥98% entity preservation', 'level': 1},
])

# Slide 9: Signal B Overview
add_content_slide(prs, "Signal B: Entity Preservation", [
    {'text': 'Core Insight: Critical information must survive pipeline'},
    {'text': 'Document intelligence extracts structured data:'},
    {'text': 'Names: "Dr. Jane Smith", "Einstein"', 'level': 1, 'size': 16},
    {'text': 'Numbers: "42", "3.14159", "$1,000,000"', 'level': 1, 'size': 16},
    {'text': 'Dates: "January 5, 1990", "2025-06-08"', 'level': 1, 'size': 16},
    {'text': 'Citations: "[Smith et al., 2023]"', 'level': 1, 'size': 16},
    {'text': 'Test: Do these entities survive cropping, binarization, PDF generation?', 'size': 20},
])

# Slide 10: Signal B - How It Works
add_diagram_slide(prs, "Entity Preservation: How It Works", [
    {'text': '1. Extract Entities from Source Image', 'size': 18, 'bold': True, 'bg_color': (230, 240, 255)},
    {'text': '   • OCR the original keyframe', 'size': 14},
    {'text': '   • NER (Named Entity Recognition): extract names, dates, numbers', 'size': 14},
    {'text': '   • Store as ground truth: GT = {entities}', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '2. Process Through Pipeline', 'size': 18, 'bold': True, 'bg_color': (230, 255, 230)},
    {'text': '   • Crop → Split → Binarize → PDF', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '3. Extract Entities from Final PDF', 'size': 18, 'bold': True, 'bg_color': (255, 240, 230)},
    {'text': '   • OCR the PDF page', 'size': 14},
    {'text': '   • NER: extract entities: FINAL = {entities}', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '4. Compare: GT ∩ FINAL / GT', 'size': 18, 'bold': True, 'bg_color': (255, 230, 255)},
])

# Slide 11: Signal B - Code Example
code_b = """def test_entity_preservation():
    # Original keyframe
    original_img = cv2.imread("keyframe_042.jpg")

    # Extract ground truth entities
    gt_text = ocr(original_img)
    gt_entities = extract_entities(gt_text)
    # gt_entities = {
    #   'names': ['Dr. Smith', 'Einstein'],
    #   'dates': ['2025-06-08'],
    #   'numbers': ['42', '3.14159']
    # }

    # Run through pipeline
    final_pdf = full_pipeline(original_img)

    # Extract entities from final PDF
    final_text = ocr_pdf(final_pdf, page=1)
    final_entities = extract_entities(final_text)

    # Self-verifying: entities must be preserved
    preservation_rate = entity_overlap(gt_entities, final_entities)
    assert preservation_rate > 0.98
    ✅ AUTOMATIC VERIFICATION"""

add_code_slide(prs, "Signal B: Implementation",
    code_b,
    "Critical entities must survive all processing steps"
)

# Slide 12: Signal B - What Counts as Entity?
add_content_slide(prs, "Signal B: Entity Types", [
    {'text': 'Names (PERSON entities)', 'size': 20},
    {'text': 'Proper nouns, author names, historical figures', 'level': 1},
    {'text': 'Numbers (QUANTITY entities)', 'size': 20},
    {'text': 'Integers, decimals, scientific notation, monetary values', 'level': 1},
    {'text': 'Dates (DATE/TIME entities)', 'size': 20},
    {'text': 'Any date format: "June 8, 2025", "2025-06-08", "08/06/25"', 'level': 1},
    {'text': 'Citations (REFERENCE entities)', 'size': 20},
    {'text': 'Academic references: "[Smith, 2023]", "Jones et al."', 'level': 1},
])

# Slide 13: Signal C Overview
add_content_slide(prs, "Signal C: Retrieval Correctness", [
    {'text': 'Core Insight: The database must answer known queries'},
    {'text': 'Document intelligence ultimate goal: searchable database'},
    {'text': 'Test Strategy:'},
    {'text': '1. Create test documents with known facts', 'level': 1},
    {'text': '2. Process through pipeline → database', 'level': 1},
    {'text': '3. Query the database with known questions', 'level': 1},
    {'text': '4. Verify answers match expected results', 'level': 1},
    {'text': 'If retrieval works, the entire pipeline worked!', 'size': 20},
])

# Slide 14: Signal C - How It Works
add_diagram_slide(prs, "Retrieval Correctness: How It Works", [
    {'text': '1. Create Test Corpus with Known Facts', 'size': 18, 'bold': True, 'bg_color': (230, 240, 255)},
    {'text': '   • Page 5: "Einstein published relativity in 1905"', 'size': 14},
    {'text': '   • Page 12: "The speed of light is 299,792,458 m/s"', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '2. Process Video → PDF → Database', 'size': 18, 'bold': True, 'bg_color': (230, 255, 230)},
    {'text': '   • Full ScanStudio pipeline', 'size': 14},
    {'text': '   • Index text in vector/document database', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '3. Query with Known Questions', 'size': 18, 'bold': True, 'bg_color': (255, 240, 230)},
    {'text': '   • Q: "When did Einstein publish relativity?"', 'size': 14},
    {'text': '   • Q: "What is the speed of light?"', 'size': 14},
    {'text': '', 'height': 0.3},
    {'text': '4. Verify Answers', 'size': 18, 'bold': True, 'bg_color': (255, 230, 255)},
    {'text': '   • Expected: "1905" → Retrieved: "1905" ✅', 'size': 14},
])

# Slide 15: Signal C - Code Example
code_c = """def test_retrieval_correctness():
    # Known test corpus: video of book with known facts
    test_video = "fixtures/einstein_biography.mp4"

    # Ground truth queries and answers
    ground_truth = {
        "When did Einstein publish relativity?": "1905",
        "What is the speed of light?": "299,792,458 m/s",
        "Who developed quantum mechanics?": "Heisenberg"
    }

    # Process through full pipeline
    pdf = scanstudio_pipeline(test_video)
    database = index_pdf_to_database(pdf)

    # Query and verify
    for question, expected_answer in ground_truth.items():
        retrieved_answer = database.query(question)

        # Check if expected answer is in retrieved result
        assert expected_answer in retrieved_answer, \\
            f"Query '{question}': expected '{expected_answer}', \\
             got '{retrieved_answer}'"

    ✅ AUTOMATIC VERIFICATION"""

add_code_slide(prs, "Signal C: Implementation",
    code_c,
    "End-to-end verification: can we retrieve known facts?"
)

# Slide 16: Signal C - Retrieval Metrics
add_content_slide(prs, "Signal C: Retrieval Metrics", [
    {'text': 'Exact Match'},
    {'text': 'Query answer exactly matches expected', 'level': 1},
    {'text': 'Threshold: 100% for factual queries', 'level': 1},
    {'text': 'Partial Match'},
    {'text': 'Answer contains expected substring', 'level': 1},
    {'text': 'Threshold: ≥95% for complex answers', 'level': 1},
    {'text': 'Retrieval Recall'},
    {'text': 'Fraction of known facts successfully retrieved', 'level': 1},
    {'text': 'Threshold: ≥90% recall across test corpus', 'level': 1},
])

# Slide 17: Comparison of Three Signals
add_content_slide(prs, "Comparing the Three Signals", [
    {'text': 'Signal A: Duplicate-Frame Agreement', 'size': 18},
    {'text': 'Tests: Pipeline consistency and stability', 'level': 1, 'size': 14},
    {'text': 'Catches: Non-deterministic bugs, quality variance', 'level': 1, 'size': 14},
    {'text': 'Signal B: Entity Preservation', 'size': 18},
    {'text': 'Tests: Information fidelity through pipeline', 'level': 1, 'size': 14},
    {'text': 'Catches: Cropping too aggressive, OCR degradation', 'level': 1, 'size': 14},
    {'text': 'Signal C: Retrieval Correctness', 'size': 18},
    {'text': 'Tests: End-to-end functional correctness', 'level': 1, 'size': 14},
    {'text': 'Catches: Any bug that breaks information access', 'level': 1, 'size': 14},
])

# Slide 18: Test Infrastructure
add_content_slide(prs, "Test Infrastructure", [
    {'text': 'Test Fixtures Required'},
    {'text': '• Test videos with intentional duplicate page captures', 'level': 1},
    {'text': '• Curated test corpus with known entities and facts', 'level': 1},
    {'text': '• Pre-labeled ground truth for 20-30 test pages', 'level': 1},
    {'text': 'Tools & Libraries'},
    {'text': '• imagehash (pHash for duplicate detection)', 'level': 1},
    {'text': '• spaCy or Hugging Face (NER for entity extraction)', 'level': 1},
    {'text': '• python-Levenshtein (text similarity)', 'level': 1},
    {'text': '• Vector DB (Pinecone/Weaviate) or keyword search', 'level': 1},
])

# Slide 19: Implementation Roadmap
add_content_slide(prs, "Implementation Roadmap", [
    {'text': 'Week 1: Test Fixtures', 'size': 20},
    {'text': 'Create test videos with duplicate captures', 'level': 1},
    {'text': 'Label ground truth entities for 20 test pages', 'level': 1},
    {'text': 'Week 2: Signal A (Duplicate-Frame)', 'size': 20},
    {'text': 'Implement duplicate detection and agreement scoring', 'level': 1},
    {'text': 'Week 3: Signal B (Entity Preservation)', 'size': 20},
    {'text': 'Integrate NER, build entity comparison tests', 'level': 1},
    {'text': 'Week 4: Signal C (Retrieval)', 'size': 20},
    {'text': 'Create test corpus, implement Q&A tests', 'level': 1},
])

# Slide 20: Success Metrics
add_content_slide(prs, "Success Metrics", [
    {'text': 'Signal A: Duplicate-Frame Agreement'},
    {'text': '≥95% text similarity across duplicates', 'level': 1},
    {'text': '≥90% bounding box IoU', 'level': 1},
    {'text': 'Signal B: Entity Preservation'},
    {'text': '≥98% entity preservation rate', 'level': 1},
    {'text': 'Zero critical entity loss (names, dates)', 'level': 1},
    {'text': 'Signal C: Retrieval Correctness'},
    {'text': '≥90% retrieval recall on test corpus', 'level': 1},
    {'text': '100% exact match for factual queries', 'level': 1},
])

# Slide 21: Benefits
add_content_slide(prs, "Benefits of This Approach", [
    {'text': '✅ Truly Self-Verifiable'},
    {'text': 'All three signals produce automatic pass/fail', 'level': 1},
    {'text': '✅ End-to-End Coverage'},
    {'text': 'Tests the complete document intelligence lifecycle', 'level': 1},
    {'text': '✅ Aligned with Research Goals'},
    {'text': 'Directly supports SIGMOD paper on ground truth maintenance', 'level': 1},
    {'text': '✅ Catches Real Bugs'},
    {'text': 'Detects issues that matter: lost information, poor retrieval', 'level': 1},
])

# Slide 22: Example Test Output
add_content_slide(prs, "What Test Results Look Like", [
    {'text': '$ pytest tests/test_three_signals.py', 'size': 14},
    {'text': '', 'size': 12},
    {'text': 'test_duplicate_frame_agreement ... ✅ PASSED', 'size': 14},
    {'text': '  Mean agreement: 96.3% (threshold: 95%)', 'level': 1, 'size': 12},
    {'text': '', 'size': 12},
    {'text': 'test_entity_preservation ... ❌ FAILED', 'size': 14},
    {'text': '  Entity loss: 3 dates missing after binarization', 'level': 1, 'size': 12},
    {'text': '  Preservation rate: 94.2% (threshold: 98%)', 'level': 1, 'size': 12},
    {'text': '', 'size': 12},
    {'text': 'test_retrieval_correctness ... ✅ PASSED', 'size': 14},
    {'text': '  Retrieved: 18/20 known facts (90%)', 'level': 1, 'size': 12},
])

# Slide 23: Next Steps
add_content_slide(prs, "Next Steps", [
    {'text': '1. Collect Test Data'},
    {'text': 'Record test videos with duplicate page captures', 'level': 1},
    {'text': 'Curate 20 pages with rich entities (names, dates, citations)', 'level': 1},
    {'text': '2. Build Test Infrastructure'},
    {'text': 'Install dependencies: imagehash, spaCy, python-Levenshtein', 'level': 1},
    {'text': 'Create test harness for three signals', 'level': 1},
    {'text': '3. Implement Tests'},
    {'text': 'Start with Signal A (easiest to implement)', 'level': 1},
    {'text': '4. Integrate into CI/CD'},
    {'text': 'Run on every commit to catch regressions early', 'level': 1},
])

# Slide 24: Questions
add_title_slide(prs,
    "Questions?",
    "Ready to implement these three verification signals"
)

# Save presentation
output_path = "ScanStudio_Testing_Strategy.pptx"
prs.save(output_path)
print(f"✅ Presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📍 Location: {output_path}")
